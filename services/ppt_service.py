# -*- coding: utf-8 -*-
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from config import PPT_DIR


class PPTService:
    """PPT 生成服务。

    将调研报告大纲自动填入 PPT 模板，生成可直接演示的幻灯片。
    支持默认模板和自定义模板。
    """

    OUTPUT_DIR = PPT_DIR

    @classmethod
    def create_pptx(
        cls,
        title: str,
        slides_data: list[dict],
        template_path: str = "",
    ) -> str:
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = cls._build_default_template()

        for slide_data in slides_data:
            slide_type = slide_data.get("type", "content")
            slide_title = slide_data.get("title", "")
            slide_content = slide_data.get("content", "")

            if slide_type == "title":
                cls._add_title_slide(prs, title, slide_content)
            elif slide_type == "section":
                cls._add_section_slide(prs, slide_title)
            else:
                cls._add_content_slide(prs, slide_title, slide_content)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c if c.isalnum() or c in "._- " else "_" for c in title)
        safe_title = safe_title.strip().replace(" ", "_")[:40]
        filename = f"ppt_{timestamp}_{safe_title}.pptx"
        filepath = os.path.join(cls.OUTPUT_DIR, filename)
        prs.save(filepath)
        return filepath

    @classmethod
    def _build_default_template(cls) -> Presentation:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        return prs

    @classmethod
    def _add_title_slide(cls, prs, title, subtitle):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(10.3)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.LEFT

        top2 = Inches(3.8)
        height2 = Inches(1.2)
        txBox2 = slide.shapes.add_textbox(left, top2, width, height2)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        p2.alignment = PP_ALIGN.LEFT

        line = slide.shapes.add_shape(
            1, left, Inches(5.3), Inches(2), Pt(3)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x1A, 0x73, 0xE8)
        line.line.fill.background()

    @classmethod
    def _add_section_slide(cls, prs, title):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)

        left = Inches(1.5)
        top = Inches(2.8)
        width = Inches(10.3)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x73, 0xE8)
        p.alignment = PP_ALIGN.LEFT

    @classmethod
    def _add_content_slide(cls, prs, title, content):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        left = Inches(1.0)
        top = Inches(0.5)
        width = Inches(11.3)
        height = Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        p.alignment = PP_ALIGN.LEFT

        line = slide.shapes.add_shape(
            1, left, Inches(1.4), width, Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        line.line.fill.background()

        top2 = Inches(1.7)
        height2 = Inches(5.2)
        txBox2 = slide.shapes.add_textbox(left, top2, width, height2)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        lines = content.strip().split("\n")
        for i, line_text in enumerate(lines):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = line_text
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(6)

    @classmethod
    def build_outline_slides(cls, outline_text: str) -> list[dict]:
        slides = []
        lines = outline_text.strip().split("\n")
        current_section = ""
        current_content: list[str] = []
        is_first = True

        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue

            if stripped.startswith("## ") or stripped.startswith("# "):
                if current_section:
                    slides.append({
                        "type": "content",
                        "title": current_section.lstrip("# ").strip(),
                        "content": "\n".join(current_content),
                    })
                current_section = stripped.lstrip("# ").strip()
                current_content = []
                if is_first:
                    is_first = False
            else:
                current_content.append(stripped)

        if current_section:
            slides.append({
                "type": "content",
                "title": current_section,
                "content": "\n".join(current_content),
            })

        return slides
